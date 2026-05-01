"""
Generate NYC Commercial Intelligence presentation (~5 min / 7 slides).
Run: python3 make_ppt.py
Output: nyc_commercial_intelligence.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
TEAL      = RGBColor(0x1A, 0x93, 0x6F)   # accent / headings
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xD4, 0xE9, 0xE2)   # soft teal for sub-text
YELLOW    = RGBColor(0xF5, 0xC5, 0x18)   # highlight boxes
GRAY_BG   = RGBColor(0x16, 0x27, 0x38)   # card background
MID_GRAY  = RGBColor(0x8A, 0xA8, 0xBE)   # secondary text

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h, *,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, title, bullets, x, y, w, h, *,
                   title_size=20, bullet_size=16,
                   bg=GRAY_BG, title_color=TEAL, bullet_color=WHITE):
    """Draws a card with a colored title and bullet list."""
    add_rect(slide, x, y, w, h, bg)

    # Title bar inside card
    bar_h = Inches(0.45)
    add_rect(slide, x, y, w, bar_h, TEAL)
    add_textbox(slide, title, x + Inches(0.15), y + Inches(0.04),
                w - Inches(0.3), bar_h,
                font_size=title_size, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Bullets
    txb = slide.shapes.add_textbox(x + Inches(0.2), y + bar_h + Inches(0.08),
                                    w - Inches(0.4), h - bar_h - Inches(0.12))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size  = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.space_before = Pt(3)

def slide_bg(slide):
    add_rect(slide, 0, 0, W, H, NAVY)

def slide_header(slide, label, title, subtitle=None):
    """Top accent bar + slide number label + big title."""
    # Accent bar
    add_rect(slide, 0, 0, W, Inches(0.07), TEAL)
    # Left label tag
    tag_w = Inches(2.8)
    add_rect(slide, 0, Inches(0.07), tag_w, Inches(0.55), TEAL)
    add_textbox(slide, label, Inches(0.18), Inches(0.07), tag_w - Inches(0.2),
                Inches(0.55), font_size=13, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    # Title
    add_textbox(slide, title,
                Inches(0.35), Inches(0.75), W - Inches(0.7), Inches(1.0),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(1.65), W - Inches(0.7), Inches(0.5),
                    font_size=18, color=LIGHT, align=PP_ALIGN.LEFT)

def code_box(slide, code_text, x, y, w, h):
    add_rect(slide, x, y, w, h, RGBColor(0x0A, 0x14, 0x1E))
    add_textbox(slide, code_text, x + Inches(0.18), y + Inches(0.1),
                w - Inches(0.3), h - Inches(0.2),
                font_size=14, color=YELLOW, align=PP_ALIGN.LEFT)

def highlight_box(slide, text, x, y, w, h, bg=YELLOW, fg=NAVY, size=16):
    add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, text, x + Inches(0.1), y + Inches(0.05),
                w - Inches(0.2), h - Inches(0.1),
                font_size=size, bold=True, color=fg, align=PP_ALIGN.CENTER, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "01  PROBLEM STATEMENT", "Why Is Commercial Location\nSelection Hard in NYC?")

# Three pain-point cards
cards = [
    ("Scale",        ["71 distinct neighborhoods", "50+ measurable signals per area", "No single 'right' answer"]),
    ("Complexity",   ["Foot traffic ≠ business success", "High density = high competition", "Demographics ≠ demand"]),
    ("Gap",          ["Existing tools: raw stats only", "No semantic / intent matching", "No competition-aware ranking"]),
]
cx = Inches(0.4)
for title, bullets in cards:
    add_bullet_box(s, title, bullets, cx, Inches(2.4), Inches(3.9), Inches(2.2),
                   title_size=18, bullet_size=15)
    cx += Inches(4.1)

# Bottom insight
highlight_box(s,
    "Goal: Blend deterministic filters + semantic search + competition scoring into one transparent ranking system",
    Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.75), size=15)

# Bottom footnote
add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Data Processing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "02  DATA PROCESSING", "Five Raw Sources → Three Clean CSVs",
             subtitle="All inputs are public NYC Open Data / MTA / DOT datasets")

sources = [
    ("Pedestrian Traffic",   "NYC DOT Bi-Annual Counts",       "→  ped_clean.csv"),
    ("Subway Stations",      "MTA Station Locations",           "→  subway_clean.csv"),
    ("Neighborhood Profiles","NYC Comptroller MOCEJ",           "→  nbhd_clean.csv"),
    ("Storefront Filings",   "NYC Open Data (optional)",        "→  storefront_features.csv"),
    ("Shooting Incidents",   "NYPD 2024 (optional)",            "→  shooting_features.csv"),
]

row_y = Inches(2.3)
for src, detail, output in sources:
    add_rect(s, Inches(0.4), row_y, Inches(4.5), Inches(0.52), GRAY_BG)
    add_textbox(s, src,    Inches(0.55), row_y + Inches(0.06), Inches(2.3), Inches(0.42),
                font_size=14, bold=True, color=WHITE)
    add_textbox(s, detail, Inches(2.85), row_y + Inches(0.06), Inches(2.0), Inches(0.42),
                font_size=12, color=MID_GRAY)
    # Arrow
    add_textbox(s, "▶", Inches(5.05), row_y + Inches(0.1), Inches(0.3), Inches(0.4),
                font_size=16, color=TEAL, align=PP_ALIGN.CENTER)
    # Output box
    add_rect(s, Inches(5.4), row_y, Inches(2.6), Inches(0.52), RGBColor(0x12, 0x3A, 0x2B))
    add_textbox(s, output, Inches(5.55), row_y + Inches(0.08), Inches(2.4), Inches(0.4),
                font_size=13, color=TEAL, bold=True)
    row_y += Inches(0.62)

# Key cleaning steps card on right
add_bullet_box(s, "Key Cleaning Steps", [
    "Standardize borough & CD keys",
    "Deduplicate station / sensor rows",
    "Normalize column names",
    "Handle optional datasets gracefully",
    "Output: zero NaN guarantee after pipeline",
], Inches(8.3), Inches(2.3), Inches(4.6), Inches(2.9), bullet_size=14)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Feature Engineering
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "03  FEATURE ENGINEERING",
             "Spatial Joins + Derived Scores",
             subtitle="Output: neighborhood_features_final.csv  ·  71 rows × 75 columns")

# Left: pipeline flow
flow = [
    "CDTA 2020 Boundary Polygons (shapefile)",
    "↓  Geopandas spatial join",
    "Aggregate pedestrian / subway / storefront\ncounts per CDTA polygon",
    "↓  Merge neighborhood profiles (CD key)",
    "Borough-median imputation for missing values",
    "↓  Derive interaction scores",
    "neighborhood_features_final.csv",
]
fy = Inches(2.3)
for i, line in enumerate(flow):
    is_step = line.startswith("↓")
    is_out  = "final" in line
    color   = TEAL if is_step else (YELLOW if is_out else WHITE)
    bold    = is_out
    size    = 13 if is_step else (15 if is_out else 14)
    add_textbox(s, line, Inches(0.4), fy, Inches(5.8), Inches(0.55),
                font_size=size, color=color, bold=bold)
    fy += Inches(0.55)

# Right: key formulas
add_bullet_box(s, "Key Derived Features", [
    "",
], Inches(6.6), Inches(2.3), Inches(6.3), Inches(0.45))  # title-only header

code_box(s,
    "commercial_activity_score\n  = log1p( storefront_filings × avg_pedestrian )\n\n"
    "competitive_score\n  = log1p( storefront_filings / avg_pedestrian )\n\n"
    "transit_activity_score\n  = log1p( subway_stations × avg_pedestrian )\n\n"
    "category_entropy  =  Shannon entropy( act_* counts )",
    Inches(6.6), Inches(2.75), Inches(6.3), Inches(2.5))

add_bullet_box(s, "Why log1p?", [
    "Compresses heavy tails (Midtown vs. park)",
    "Keeps 0 as 0 (no log-of-zero error)",
    "Makes slider thresholds intuitive",
], Inches(6.6), Inches(5.4), Inches(6.3), Inches(1.65), bullet_size=13)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Clustering
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "04  CLUSTERING",
             "K-Means From Scratch (NumPy)",
             subtitle="Exploratory tool for discovering neighborhood archetypes")

add_bullet_box(s, "Implementation  (src/kmeans_numpy.py)", [
    "Pure NumPy — no scikit-learn for core algorithm",
    "Z-score normalization before clustering",
    "Euclidean distance, iterative centroid update",
    "Convergence: centroid shift < 1e-4  OR  100 iterations",
    "Empty cluster reinit from random points",
    "HDF5 cache per k value",
], Inches(0.4), Inches(2.3), Inches(6.0), Inches(2.9), bullet_size=14)

add_bullet_box(s, "K Selection", [
    "Sweep k = 2 … max_k",
    "Elbow curve  (inertia / WCSS)",
    "Silhouette score — also from scratch",
    "Auto-detect elbow: perpendicular-distance method",
    "Optional: cluster briefs via embeddings",
], Inches(0.4), Inches(5.35), Inches(6.0), Inches(1.8), bullet_size=14)

# Right: features used
add_bullet_box(s, "Default Clustering Features", [
    "storefront_filing_count",
    "avg_pedestrian",
    "subway_station_count",
    "storefront_density_per_km2",
    "commercial_activity_score",
    "competitive_score",
    "shooting_incident_count",
    "category_entropy",
    "nfh_overall_score  (if available)",
], Inches(6.6), Inches(2.3), Inches(6.3), Inches(3.3), bullet_size=13)

highlight_box(s,
    "Purpose: discover archetypes (e.g. 'high-traffic commercial', 'residential low-competition')\n"
    "Not used for ranking — purely exploratory on the home page",
    Inches(6.6), Inches(5.75), Inches(6.3), Inches(0.95), size=13)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Embeddings + Constraints
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "05  EMBEDDINGS + CONSTRAINTS",
             "Text Profiles → Vectors + Filter Layer",
             subtitle="Bridging structured data and natural-language business intent")

# Left: text profile pipeline
add_bullet_box(s, "Text Profile Generation  (build_text_profile)", [
    "Every neighborhood row → natural language paragraph",
    "Numbers mapped to labels:  avg_ped > 5000 → 'very high foot traffic'",
    "All act_* counts written as:  'FOOD SERVICES (3789), RETAIL (6995)…'",
    "Demographics, jobs, NFH scores, crime all included",
    "→  ~300 word profile per neighborhood",
], Inches(0.4), Inches(2.3), Inches(6.2), Inches(2.6), bullet_size=13)

code_box(s,
    "# Example snippet (Midtown-Flatiron)\n"
    "\"19752 storefront filings; extremely dense (5391/km2).\n"
    " very high foot traffic (avg 8363 pedestrians).\n"
    " Counts: RETAIL (6995), FOOD SERVICES (3789),\n"
    "         EDUCATIONAL SERVICES (2717)...\"",
    Inches(0.4), Inches(5.05), Inches(6.2), Inches(2.1))

# Right: constraints
add_bullet_box(s, "Hard Constraints  (DuckDB SQL)", [
    "Borough, min subway stations, min pedestrians",
    "Min storefront density / filing count",
    "Max competitive score  (market saturation cap)",
    "Max shooting incidents  (safety floor)",
    "Min NFH financial-health scores  (optional)",
    "→  Deterministic, transparent SQL output",
], Inches(6.8), Inches(2.3), Inches(6.1), Inches(2.6), bullet_size=13)

add_bullet_box(s, "Soft Constraints  (Semantic)", [
    "User types free-text: 'quiet boutique retail area'",
    "Query embedded with same model as profiles",
    "Cosine similarity → semantic_similarity ∈ [0, 1]",
    "Only applied to hard-filtered candidates",
], Inches(6.8), Inches(5.05), Inches(6.1), Inches(2.1), bullet_size=13)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Ranking
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "06  RANKING",
             "Two-Stage Blended Scoring",
             subtitle="Deterministic filters first, then weighted semantic + competition blend")

# Stage diagram
add_rect(s, Inches(0.4), Inches(2.2), Inches(2.5), Inches(0.65), TEAL)
add_textbox(s, "All 71 Neighborhoods", Inches(0.5), Inches(2.28), Inches(2.35), Inches(0.5),
            font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(s, "▼  Stage 1: DuckDB SQL Hard Filters",
            Inches(0.4), Inches(3.05), Inches(3.5), Inches(0.45),
            font_size=13, color=TEAL)

add_rect(s, Inches(0.4), Inches(3.55), Inches(2.5), Inches(0.65), GRAY_BG)
add_textbox(s, "Filtered Candidates (n ≤ 71)", Inches(0.5), Inches(3.63), Inches(2.35), Inches(0.5),
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, "▼  Stage 2: Soft Blending",
            Inches(0.4), Inches(4.35), Inches(3.5), Inches(0.45),
            font_size=13, color=TEAL)

add_rect(s, Inches(0.4), Inches(4.85), Inches(2.5), Inches(0.65), RGBColor(0x12, 0x3A, 0x2B))
add_textbox(s, "Ranked Results ↑", Inches(0.5), Inches(4.93), Inches(2.35), Inches(0.5),
            font_size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Blend formula
code_box(s,
    "# Blended score formula\n\n"
    "specific_competitive = log1p( category_count / avg_pedestrian )\n\n"
    "[semantic, -competitive]  →  MinMaxScaler( filtered_rows )\n\n"
    "blended_score = α × semantic_similarity\n"
    "              + (1-α) × (1 - competitive_normalized)",
    Inches(3.3), Inches(2.2), Inches(5.8), Inches(3.4))

# Alpha legend
add_bullet_box(s, "α Slider Meaning", [
    "α = 0.0  →  pure competition avoidance",
    "α = 0.5  →  balanced (default)",
    "α = 1.0  →  pure semantic match",
    "MinMax fitted on filtered set only",
], Inches(9.3), Inches(2.2), Inches(3.65), Inches(2.2), bullet_size=13)

highlight_box(s,
    "Why MinMax on filtered set?\nEnsures α=0.5 is truly balanced regardless of score scale",
    Inches(9.3), Inches(4.55), Inches(3.65), Inches(0.95), size=13)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Validation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
slide_header(s, "07  VALIDATION",
             "Tests, Rank Stability & Live Demo",
             subtitle="Correctness at the algorithm level + behavioral sanity at the system level")

add_bullet_box(s, "Unit Tests  (pytest  tests/)", [
    "K-means: convergence, label assignment, inertia, silhouette",
    "K-means caching: HDF5 save → load round-trip",
    "Feature engineering: spatial join, aggregation correctness",
    "MiniBatch k-means: NotImplementedError guard",
], Inches(0.4), Inches(2.3), Inches(6.2), Inches(2.3), bullet_size=14)

add_bullet_box(s, "Rank Stability Validation", [
    "rank_stability_validation_business_queries.py",
    "Fixed set of business queries (e.g. 'quiet boutique retail')",
    "Assert top-3 neighborhoods are stable across α values",
    "Catch regressions when scoring formula changes",
], Inches(0.4), Inches(4.75), Inches(6.2), Inches(2.4), bullet_size=14)

add_bullet_box(s, "System-Level Checks", [
    "Zero NaN guarantee logged after pipeline run",
    "Embedding cache validated on load (shape check)",
    "DuckDB SQL generated and shown to user (transparent)",
    "Claude agent SQL: read-only SELECT enforced",
], Inches(6.8), Inches(2.3), Inches(6.1), Inches(2.3), bullet_size=14)

highlight_box(s,
    "Live Demo\n\n"
    "Query: 'high foot traffic area good for Asian restaurant'\n"
    "→  Hard filter: Manhattan, subway ≥ 5, pedestrian ≥ 3000\n"
    "→  Soft: α = 0.6, category = FOOD_SERVICES\n"
    "→  Top result: Flushing-Murray Hill  (high Asian pop + dense food scene)",
    Inches(6.8), Inches(4.75), Inches(6.1), Inches(2.4), bg=GRAY_BG, fg=WHITE, size=13)

add_textbox(s, "NYC Commercial Intelligence  ·  2025",
            Inches(0.35), Inches(6.9), Inches(6), Inches(0.4),
            font_size=11, color=MID_GRAY)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "nyc_commercial_intelligence.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
